"""
Generate fail PDF/Word/PowerPoint/Excel dari teks (biasanya jawapan terakhir AI) untuk
dihantar balik kepada student — supaya senang simpan/print/share nota.

Ni PERCUMA (tak guna AI generation) — cuma convert teks yang AI dah jana kepada
format fail, guna library biasa (fpdf2, python-docx, python-pptx, openpyxl).
"""

import io
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.styles import Font


def _sanitize_for_pdf(text: str) -> str:
    """fpdf2 (guna font core Helvetica) tak sokong semua unicode (contoh emoji).
    Tukar/buang character yang tak disokong supaya tak crash masa generate PDF."""
    return text.encode("latin-1", errors="replace").decode("latin-1")


def text_to_pdf_bytes(title: str, body: str) -> bytes:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, _sanitize_for_pdf(title))
    pdf.ln(4)
    pdf.set_font("Helvetica", size=12)
    pdf.multi_cell(0, 8, _sanitize_for_pdf(body))
    return bytes(pdf.output())


def text_to_docx_bytes(title: str, body: str) -> bytes:
    doc = Document()
    doc.add_heading(title, level=1)
    for para in body.split("\n\n"):
        if para.strip():
            doc.add_paragraph(para.strip())
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def text_to_pptx_bytes(title: str, body: str, max_chars_per_slide: int = 500) -> bytes:
    """Pecahkan teks jadi beberapa slide — slide tajuk + slide kandungan (bullet points)."""
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Dijana oleh StudyBot"

    paragraphs = [p.strip() for p in body.split("\n\n") if p.strip()] or [body.strip() or "Tiada kandungan"]

    bullet_layout = prs.slide_layouts[1]  # Title and Content
    slide_num = 1
    current_group: list[str] = []
    current_len = 0

    def flush_slide(paras: list[str]):
        nonlocal slide_num
        if not paras:
            return
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = f"Nota — bahagian {slide_num}"
        tf = slide.placeholders[1].text_frame
        tf.text = paras[0][:300]
        for p in paras[1:]:
            tf.add_paragraph().text = p[:300]
        slide_num += 1

    for para in paragraphs:
        current_group.append(para)
        current_len += len(para)
        if current_len > max_chars_per_slide or len(current_group) >= 4:
            flush_slide(current_group)
            current_group, current_len = [], 0

    flush_slide(current_group)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def text_to_xlsx_bytes(title: str, body: str) -> bytes:
    """Pecahkan teks ikut baris jadi jadual senang dibaca — sesuai untuk senarai/nota ringkas."""
    wb = Workbook()
    ws = wb.active
    ws.title = "StudyBot Nota"

    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14)
    ws.merge_cells("A1:B1")

    ws["A3"] = "Bil"
    ws["B3"] = "Kandungan"
    ws["A3"].font = Font(bold=True)
    ws["B3"].font = Font(bold=True)

    lines = [line.strip().lstrip("-•* ").strip() for line in body.split("\n") if line.strip()]
    row = 4
    for i, line in enumerate(lines, start=1):
        ws.cell(row=row, column=1, value=i)
        ws.cell(row=row, column=2, value=line)
        row += 1

    ws.column_dimensions["A"].width = 6
    ws.column_dimensions["B"].width = 90

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
